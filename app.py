import streamlit as st
import os
import tempfile
from langchain_community.document_loaders import PyPDFLoader
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import (
    GoogleGenerativeAIEmbeddings,
    ChatGoogleGenerativeAI
)
from langchain.vectorstores import FAISS
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from pptx import Presentation

# --------------------------------------------------
# PAGE CONFIG
# --------------------------------------------------
st.set_page_config(
    page_title="AI Exam Tutor (RAG + Gemini)",
    page_icon="📘",
    layout="wide"
)

st.title("📘 AI Exam Tutor (PDF & PPT Based)")
st.caption("Answers strictly from uploaded documents • Simple explanations")

# --------------------------------------------------
# API KEY
# --------------------------------------------------
if "GOOGLE_API_KEY" not in st.secrets:
    st.error("❌ Add GOOGLE_API_KEY in Streamlit Secrets")
    st.stop()

os.environ["GOOGLE_API_KEY"] = st.secrets["GOOGLE_API_KEY"]

# --------------------------------------------------
# SESSION STATE
# --------------------------------------------------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# --------------------------------------------------
# FILE UPLOAD
# --------------------------------------------------
uploaded_files = st.file_uploader(
    "Upload PDF or PPT files",
    type=["pdf", "pptx"],
    accept_multiple_files=True
)

documents = []

# --------------------------------------------------
# LOADERS
# --------------------------------------------------
def load_pdf(file):
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(file.read())
        loader = PyPDFLoader(tmp.name)
        return loader.load()

def load_ppt(file):
    prs = Presentation(file)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(f"Slide {i+1}: {shape.text}")
    return [Document(page_content="\n".join(slides_text), metadata={"source": file.name})]

# --------------------------------------------------
# PROCESS FILES
# --------------------------------------------------
if uploaded_files:
    for file in uploaded_files:
        if file.name.endswith(".pdf"):
            docs = load_pdf(file)
            for d in docs:
                d.metadata["source"] = file.name
            documents.extend(docs)
        elif file.name.endswith(".pptx"):
            documents.extend(load_ppt(file))

# --------------------------------------------------
# BUILD VECTOR STORE
# --------------------------------------------------
if documents:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100
    )
    chunks = splitter.split_documents(documents)

    embeddings = GoogleGenerativeAIEmbeddings(
        model="models/embedding-001"
    )

    vectorstore = FAISS.from_documents(chunks, embeddings)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})

    # --------------------------------------------------
    # USER CONTROLS
    # --------------------------------------------------
    col1, col2 = st.columns(2)
    with col1:
        answer_mode = st.selectbox(
            "Answer Style",
            ["Simple Explanation", "Bullet Points"]
        )
    with col2:
        show_sources = st.checkbox("Show source reference", value=True)

    # --------------------------------------------------
    # PROMPT
    # --------------------------------------------------
    prompt = PromptTemplate(
        input_variables=["context", "question"],
        template=f"""
You are an exam-focused AI tutor.

Rules:
- Answer ONLY from the given context
- Use very simple language
- Do NOT add outside knowledge
- If answer is missing, say:
  "Answer not found in the uploaded documents."

Answer style: {answer_mode}

Context:
{{context}}

Question:
{{question}}

Answer:
"""
    )

    llm = ChatGoogleGenerativeAI(
        model="gemini-1.5-flash",
        temperature=0.2
    )

    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        chain_type="stuff",
        chain_type_kwargs={"prompt": prompt},
        return_source_documents=True
    )

    # --------------------------------------------------
    # QUESTION INPUT
    # --------------------------------------------------
    st.divider()
    question = st.text_input("Ask a question from the uploaded files")

    if question:
        with st.spinner("Analyzing documents..."):
            result = qa_chain(question)

        answer = result["result"]
        sources = result["source_documents"]

        st.session_state.chat_history.append((question, answer))

        st.subheader("Answer")
        st.write(answer)

        if show_sources:
            st.subheader("Source Reference")
            for doc in sources:
                src = doc.metadata.get("source", "Unknown")
                page = doc.metadata.get("page", "N/A")
                st.caption(f"📄 {src} | Page/Slide: {page}")

    # --------------------------------------------------
    # CHAT HISTORY
    # --------------------------------------------------
    if st.session_state.chat_history:
        st.divider()
        st.subheader("Previous Questions")
        for i, (q, a) in enumerate(reversed(st.session_state.chat_history), 1):
            with st.expander(f"Q{i}: {q}"):
                st.write(a)

    # --------------------------------------------------
    # DOCUMENT SUMMARY
    # --------------------------------------------------
    st.divider()
    if st.button("Generate Summary of Documents"):
        with st.spinner("Generating summary..."):
            summary_prompt = """
Summarize the following content in easy exam-oriented points.
"""
            summary = llm.invoke(summary_prompt + "\n".join([c.page_content for c in chunks[:6]]))
        st.subheader("Document Summary")
        st.write(summary.content)

else:
    st.info("📂 Upload at least one PDF or PPT to begin.")
